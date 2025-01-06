import openai
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import streamlit as st
from dotenv import load_dotenv
import requests
from io import BytesIO

# Load OpenAI API Key
load_dotenv()
openai.api_key = os.getenv('OPENAI_API_KEY')

# Define formatting options
TITLE_FONT_SIZE = Pt(20)
TEXT_FONT_SIZE = Pt(14)

# Function to fetch images from Unsplash API
def fetch_image(topic):
    UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")  # Replace with your Unsplash API Key
    url = f"https://api.unsplash.com/photos/random?query={topic}&client_id={UNSPLASH_ACCESS_KEY}"
    response = requests.get(url)
    if response.status_code == 200:
        image_url = response.json().get("urls", {}).get("regular")
        return image_url
    else:
        return None

# Function to generate slide titles
def generate_slide_titles(topic):
    try:
        messages = [
            {"role": "system", "content": "You are an AI assistant that helps create presentations."},
            {"role": "user", "content": f"Generate 10 slide titles for the topic '{topic}'."}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=100,
        )
        return response['choices'][0]['message']['content'].split("\n")
    except Exception as e:
        st.error(f"Error generating slide titles: {e}")
        return []

# Function to generate slide content
def generate_slide_content(slide_title):
    try:
        messages = [
            {"role": "system", "content": "You are an AI assistant that helps create presentations."},
            {"role": "user", "content": f"Generate content for the slide: '{slide_title}'."}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=100,
        )
        return response['choices'][0]['message']['content']
    except Exception as e:
        st.error(f"Error generating slide content for '{slide_title}': {e}")
        return ""

# Function to create a presentation with a template
def create_presentation_with_template(template_path, topic, slide_titles, slide_contents):
    prs = pptx.Presentation(template_path)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]
    title_shape.text = topic
    subtitle_shape.text = "Generated with AI and Streamlit"

    # Customize title slide font and color
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = TITLE_FONT_SIZE
        paragraph.font.color.rgb = RGBColor(0, 102, 204)

    # Content slides
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide_layout = prs.slide_layouts[1]  # Assuming content layout from template
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.title
        title_box.text = slide_title
        for paragraph in title_box.text_frame.paragraphs:
            paragraph.font.size = TITLE_FONT_SIZE

        # Add content text on the left
        text_box = slide.placeholders[1]
        text_box.text = slide_content
        for paragraph in text_box.text_frame.paragraphs:
            paragraph.font.size = TEXT_FONT_SIZE

        # Fetch and add an image on the right side
        image_url = fetch_image(topic)
        if image_url:
            response = requests.get(image_url)
            if response.status_code == 200:
                img_stream = BytesIO(response.content)
                slide.shapes.add_picture(img_stream, Inches(4.5), Inches(1.5), width=Inches(4))

    # Save presentation
    output_dir = "generated_ppt"
    os.makedirs(output_dir, exist_ok=True)
    presentation_path = os.path.join(output_dir, f"{topic}_presentation.pptx")
    prs.save(presentation_path)
    return presentation_path


custom_css = """
<style>

#offer-div {
    margin-top: 30px;
}

#text-offer-div{
    text-align: center;
}

#row3-div {
    display: flex;
    justify-content: space-around;
    width: 100%;
    overflow: auto;
    padding: 10px;
    margin-bottom: 70px;
}

#chat-div, #summarize-div, #Audio-div {
    width: 30%;
    text-align: center;
    padding: 10px;
    background-color: #f2f2f2;
    border-radius: 50px;
    border-style: solid;
    border-color: #7f7f7f;
    border-width: 2px;
    border-radius: 10px;
}

#row3-div h4, #row3-div h5{
    color: black;
}


#transform-div{
        margin-top: 20px;
        position: relative;
}

#transform-div h1{
    text-align: center;
}
#transform-div h6{
    color: #afafaf;
    text-align: center;
    padding: 10px 50px;
}





</style>
"""

# Streamlit app
def main():
    st.set_page_config(layout="wide")
    
    # Apply CSS
    st.markdown(custom_css, unsafe_allow_html=True)

    # --------------------------------------------- Main section ------------------------------------------- #
    st.markdown('''
        <div id="transform-div">
            <h1>Text to PowerPoint 💻</h1>
            <h6>Welcome to the Text to PowerPoint page! Effortlessly turn your text inputs into professional PowerPoint presentations. Simply provide your content, and our tool will generate slides tailored to your needs. Perfect for students, professionals, and educators, this page saves you time while delivering impactful presentations. Start creating your slides with just a few clicks!</h6>
        </div>
    ''', unsafe_allow_html=True)

    # --------------------------------------------- How it works? section ------------------------------------------- #
    st.markdown('''
    <div id="offer-div">
        <div id="text-offer-div">
            <h1>How it works ?</h1>
        </div>
        <div id="row3-div">
            <div id="chat-div">
                <img src="https://cdn0.iconfinder.com/data/icons/education-and-learning-round/128/18-512.png" width="50" height="50">
                <h4>Enter the the presentation topic</h4>
            </div>
            <div id="summarize-div"> 
                <img src="https://static.thenounproject.com/png/2704149-200.png" width="50" height="50">
                <h5>Wait for AI to create the presentation</h5>
            </div>
            <div id="Audio-div">
                <img src="https://icon-library.com/images/icon-download/icon-download-4.jpg" width="50" height="50">
                <h4>Download your presentation</h4>
            </div>
        </div>     
    </div>

    ''', unsafe_allow_html=True)

    col1, col2 = st.columns(2)

    with col1:
        topic = st.text_input("Enter the topic for your presentation:")

    with col2:
        template_path = st.text_input("If you want to use a template Please enter the path to your PowerPoint template (.potx):")
    
    generate_button = st.button("Generate Presentation", use_container_width=True)

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles = [title for title in slide_titles if title.strip()]
        slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
        presentation_path = create_presentation_with_template(template_path, topic, filtered_slide_titles, slide_contents)

        st.success("Presentation generated successfully!")
        st.download_button(
            label="Download the PowerPoint Presentation",
            data=open(presentation_path, "rb").read(),
            file_name=f"{topic}_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

if __name__ == "__main__":
    main()
